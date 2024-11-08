from gen_text import gen_text
import json
import pptx
from pptx.util import Inches, Length
from pptx.enum.text import MSO_AUTO_SIZE

def gen_pptx(arguments: str, path: str):
    presentation = pptx.Presentation()

    presentation.slide_height = Inches(8.27)
    presentation.slide_width = Inches(11.69)

    text = gen_text(arguments)
    text = text[text.find("["):text.find("]")+1]

    slides = json.loads(text)

    for slide_obj in slides:
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)

        if slide.shapes.title is None \
            or presentation.slide_width is None \
            or presentation.slide_height is None: return

        slide.shapes.title.top = Inches(0)
        slide.shapes.title.left = Inches(0)
        slide.shapes.title.width = presentation.slide_width;
        slide.shapes.title.height = Length(presentation.slide_height / 8.0)

        slide.shapes.title.text = slide_obj["title"]

        slide.placeholders[1].top = slide.shapes.title.height + Inches(1)
        slide.placeholders[1].left = Inches(0)
        slide.placeholders[1].width = presentation.slide_width
        slide.placeholders[1].height = presentation.slide_height - slide.placeholders[1].top

        slide.placeholders[1].text = slide_obj["content"]

    presentation.save(path)

